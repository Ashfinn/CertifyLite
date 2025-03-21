from pptx import Presentation
import os
import comtypes.client  # Required for PDF conversion (Windows only)
import pandas as pd  # For reading Excel files
import time  # For adding a delay

def load_names_from_excel(excel_path):
    """ Load names from an Excel file """
    df = pd.read_excel(excel_path)
    if "NAME" not in df.columns:
        raise ValueError("Excel file must contain a 'NAME' column")
    return df["NAME"].tolist()

def generate_and_convert_certificates(template_path, names, output_folder):
    """ Generate certificates and convert directly to PDF """
    # Ensure output folder exists
    os.makedirs(output_folder, exist_ok=True)
    
    # Convert paths to absolute to avoid resolution issues
    template_path = os.path.abspath(template_path)
    output_folder = os.path.abspath(output_folder)
    
    # Initialize PowerPoint COM object
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1  # Optional: Make it visible
    
    try:
        for name in names:
            # Create a copy of the template
            new_prs = Presentation(template_path)
            
            # Replace the placeholder text
            for slide in new_prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "Arnab Aich" in run.text:  # Replace placeholder
                                    run.text = run.text.replace("Arnab Aich", name)
            
            # Sanitize filename
            safe_name = "".join(c if c.isalnum() or c in "._-" else "_" for c in name)
            pptx_path = os.path.join(output_folder, f"Certificate_{safe_name}.pptx")
            pdf_path = os.path.join(output_folder, f"Certificate_{safe_name}.pdf")
            
            # Save as PPTX temporarily
            new_prs.save(pptx_path)
            print(f"Temporary PPTX saved: {pptx_path}")
            
            # Small delay to ensure file is written
            time.sleep(0.5)
            
            # Convert to PDF
            try:
                presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
                presentation.SaveAs(pdf_path, 32)  # 32 = PDF format
                presentation.Close()
                print(f"PDF saved: {pdf_path}")
                
                # Delete the temporary PPTX file
                os.remove(pptx_path)
                print(f"Temporary PPTX deleted: {pptx_path}")
            except Exception as e:
                print(f"Error converting {pptx_path} to PDF: {e}")
    
    finally:
        # Ensure PowerPoint is closed even if an error occurs
        powerpoint.Quit()

# Example usage
excel_file = "E:/Projects/cumc-certs/names.xlsx"  # Absolute path to your Excel file
pptx_template = "E:/Projects/cumc-certs/cert.pptx"  # Absolute path to your template
output_directory = "E:/Projects/cumc-certs/Certificates"  # Absolute path to output folder

# Load names from Excel
names_list = load_names_from_excel(excel_file)

# Generate certificates and convert to PDF
generate_and_convert_certificates(pptx_template, names_list, output_directory)